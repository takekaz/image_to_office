
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from PIL import Image
import os
import io
import re

def excel_col_to_pixels(col_letter, char_width_pixels=7.0, default_col_width_chars=8.43):
    """
    Excelの列幅（文字数単位）をピクセルに変換します。
    ここでは簡略化のため、デフォルトの列幅と文字幅ピクセルを使用します。
    正確な変換はExcelのデフォルトフォント設定に依存します。
    """
    if col_letter is None:
        return 0

    col_index = 0
    for char in col_letter.upper():
        col_index = col_index * 26 + (ord(char) - ord('A') + 1)
    
    # Excelの列幅単位は1/256 of the width of the "0" character.
    # Openpyxlのデフォルトは8.43文字幅。これは約64ピクセル (96DPI時)
    # したがって、1文字幅あたり約7.59ピクセル。
    # ここでは仮に1文字あたり8ピクセルとして計算します。
    # 実際の幅は、その列のすべてのセルの最も広いコンテンツに合わせて自動調整されるため、
    # 特定の値に正確に変換するのは難しいです。
    # ここでは単純に列インデックスに文字幅を乗じています。
    return col_index * default_col_width_chars * char_width_pixels


def excel_row_to_pixels(row_index, default_row_height_pts=15.0, dpi=96):
    """
    Excelの行高（ポイント単位）をピクセルに変換します。
    1ポイント = 1/72インチ。ピクセル = インチ * DPI
    """
    # Excelの行インデックスは1から始まるため、そのまま使用
    # ここでは特定の行の高さではなく、行インデックスに対応する累積高さを計算する
    # しかし、実際のexcel_posは単一のセルであり、累積高さを計算するのではなく、
    # そのセルの左上隅の絶対位置を計算する必要がある。
    # そのため、行インデックスをそのまま使用するのではなく、指定された行の高さを使用する。
    # デフォルトの行の高さは15pt (約20ピクセル)
    return row_index * default_row_height_pts / 72 * dpi


def convert_excel_pos_to_inches(excel_pos, conversion_params):
    """
    Excelのセル座標（例: "B2"）をインチ単位のX, Y座標に変換します。
    
    Args:
        excel_pos (str): "B2"のようなExcelのセル座標。
        conversion_params (dict): 変換のためのパラメータ。
                                  例: {
                                      "col_width_pixels_per_char": 7.0, # 1文字あたりの列幅ピクセル
                                      "default_col_width_chars": 8.43, # Excelのデフォルト列幅（文字数）
                                      "row_height_points_per_row": 15.0, # 1行あたりの行高ポイント
                                      "dpi": 96 # DPI
                                  }
    Returns:
        tuple: (x_inches, y_inches)
    """
    col_letter_match = re.match(r"([A-Z]+)", excel_pos, re.IGNORECASE)
    row_num_match = re.search(r"(\d+)", excel_pos)

    if not col_letter_match or not row_num_match:
        raise ValueError(f"Invalid excel_pos format: {excel_pos}")

    col_letter = col_letter_match.group(1).upper()
    row_num = int(row_num_match.group(1))

    col_width_pixels_per_char = conversion_params.get("col_width_pixels_per_char", 7.0)
    default_col_width_chars = conversion_params.get("default_col_width_chars", 8.43)
    row_height_points_per_row = conversion_params.get("row_height_points_per_row", 15.0)
    dpi = conversion_params.get("dpi", 96)

    # X座標（列の開始位置）を計算
    x_pixels = excel_col_to_pixels(col_letter, col_width_pixels_per_char, default_col_width_chars)
    x_inches = x_pixels / dpi

    # Y座標（行の開始位置）を計算
    y_pixels = excel_row_to_pixels(row_num -1, row_height_points_per_row, dpi) # row_num-1 because row_num is 1-indexed, and we need sum of previous rows
    y_inches = y_pixels / dpi

    return x_inches, y_inches

def embed_images_to_ppt(ppt_file_path, image_folder_path, regions_and_coords, excel_conversion_params):
    """
    指定された画像フォルダの画像をPowerPointファイルに埋め込みます。
    画像内の特定の領域を抽出し、PowerPointスライドの指定座標に貼り付けます。
    スライドの右上には画像ファイル名を表記します。

    Args:
        ppt_file_path (str): 生成するPowerPointファイルのパス。既存のファイルは上書きされます。
        image_folder_path (str): 画像が保存されているフォルダのパス。
        regions_and_coords (list): 領域とセル座標のペアのリスト。
                                   例: [
                                           {"img_region": [x1, y1, x2, y2], "excel_pos": "A1"},
                                           ...
                                       ]
        excel_conversion_params (dict): Excelセル座標をインチに変換するためのパラメータ。
                                        convert_excel_pos_to_inches関数のconversion_paramsを参照。
    """
    prs = Presentation()

    # 画像フォルダ内のファイルをリストアップし、ソート
    image_files = [f for f in os.listdir(image_folder_path) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))]
    image_files.sort() # ファイル名で昇順ソート

    # 画像ごとにスライドを作成
    for filename in image_files:
        image_path = os.path.join(image_folder_path, filename)
        
        try:
            img_full = Image.open(image_path)
        except Exception as e:
            print(f"Warning: Could not open image {filename}. Skipping. Error: {e}")
            continue

        slide_layout = prs.slide_layouts[5]  # 白紙のスライドレイアウトを選択 (適宜調整)
        slide = prs.slides.add_slide(slide_layout)

        # スライド右上に画像ファイル名を追加
        title_placeholder = slide.shapes.title
        title_placeholder.text = filename
        # タイトル位置を右上へ調整 (手動で座標指定)
        title_placeholder.left = Inches(prs.slide_width.inches - 2) # 仮に右から2インチ
        title_placeholder.top = Inches(0.1) # 仮に上から0.1インチ
        title_placeholder.width = Inches(1.9)
        title_placeholder.height = Inches(0.5)
        # フォントサイズ調整
        text_frame = title_placeholder.text_frame
        text_frame.word_wrap = False
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14) # サイズ調整

        for item in regions_and_coords:
            img_region = item["img_region"]
            excel_pos = item["excel_pos"]

            try:
                # 画像の領域をクロップ
                cropped_img = img_full.crop(img_region)

                # pptxがファイルパスまたはバイナリデータから画像を読み込むため、
                # クロップした画像を一時的にメモリに保存
                img_byte_arr = io.BytesIO()
                # PIL Imageのsaveメソッドで品質を指定する場合はformatも必要
                cropped_img.save(img_byte_arr, format='PNG') 
                img_byte_arr.seek(0) # ストリームの先頭に戻る

                # Excelセル座標をインチに変換
                x_inches, y_inches = convert_excel_pos_to_inches(excel_pos, excel_conversion_params)

                # 画像をスライドに埋め込む
                # width, heightはオプション。指定しない場合は元の画像サイズで挿入される
                # ここではcropped_imgのサイズに基づいて計算
                img_width_px, img_height_px = cropped_img.size
                insert_width_inches = img_width_px / excel_conversion_params.get("dpi", 96)
                insert_height_inches = img_height_px / excel_conversion_params.get("dpi", 96)

                slide.shapes.add_picture(img_byte_arr, Inches(x_inches), Inches(y_inches),
                                         width=Inches(insert_width_inches), height=Inches(insert_height_inches))
                print(f"Embedded region {img_region} from {filename} to slide at ({x_inches:.2f}in, {y_inches:.2f}in) (Excel: {excel_pos})")

            except Exception as e:
                print(f"Warning: Could not crop or embed region {img_region} from {filename} to {excel_pos}. Error: {e}")
                continue

    # PowerPointファイルを保存
    try:
        prs.save(ppt_file_path)
        print(f"PowerPoint file saved successfully to {ppt_file_path}")
    except Exception as e:
        print(f"Error: Could not save PowerPoint file to {ppt_file_path}. Error: {e}")

if __name__ == '__main__':
    # サンプルの使用方法
    output_ppt_path = "/workspace/image_to_office/output_images.pptx"
    input_image_folder = "/workspace/image_to_office/img"
    
    # 領域とセル座標のペアの例 (ExcelのB2, C5, E8, G11セルを想定)
    sample_regions_coords = [
        {"img_region": [50, 100, 200, 150], "excel_pos": "B2"},  # 黄色文字の領域を想定
        {"img_region": [100, 200, 250, 250], "excel_pos": "C5"}, # 白文字の領域を想定
        {"img_region": [150, 300, 300, 350], "excel_pos": "E8"}, # 緑文字の領域を想定
        {"img_region": [500, 400, 650, 450], "excel_pos": "G11"} # 赤文字の領域を想定
    ]

    # Excelからインチへの変換パラメータ
    # これは非常に近似的な値であり、正確なExcelのレンダリングとは異なる場合があります。
    # 実際のExcelの列幅や行高は、フォント、DPI、列の内容によって動的に変わるため、
    # 特定のExcelファイルからの正確な変換にはopenpyxlでファイル自体を読み込む必要があります。
    # ここでは一般的なデフォルト値に基づいた仮の変換値を使用します。
    conversion_params = {
        "col_width_pixels_per_char": 7.0, # Excelデフォルトフォントでの1文字あたりの平均ピクセル幅 (目安)
        "default_col_width_chars": 8.43,  # Excelのデフォルト列幅（標準フォントの文字数単位）
        "row_height_points_per_row": 15.0, # Excelのデフォルト行高（ポイント単位）
        "dpi": 96 # 一般的な画面DPI
    }

    # 関数を実行
    embed_images_to_ppt(output_ppt_path, input_image_folder, sample_regions_coords, conversion_params)
