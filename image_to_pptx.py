


import os
import json
import tempfile # 追加
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

def excel_coord_to_inches(excel_pos: str, params: dict):
    """
    Excelのセル座標をPowerPointのインチ座標に変換します。
    この関数は仮の実装であり、実際のExcelの列幅や行高さに基づいて
    より複雑な計算が必要になる場合があります。

    Args:
        excel_pos (str): Excelのセル座標 (例: "B2")
        params (dict): 変換に必要なパラメータ (例: {"col_width_pix": 64, "row_height_pix": 20, "dpi": 96})

    Returns:
        tuple: (x_inches, y_inches) - PowerPoint上でのx, y座標 (インチ単位)
    """
    col_width_pix = params.get("col_width_pix", 64) # Excelのデフォルト列幅（ピクセル）
    row_height_pix = params.get("row_height_pix", 20) # Excelのデフォルト行高さ（ピクセル）
    dpi = params.get("dpi", 96) # スクリーンDPI

    # Excel座標を列と行のインデックスに変換
    col_str = "".join(filter(str.isalpha, excel_pos)).upper()
    row_str = "".join(filter(str.isdigit, excel_pos))

    col_idx = 0
    for char in col_str:
        col_idx = col_idx * 26 + (ord(char) - ord('A') + 1)
    row_idx = int(row_str)

    # ピクセル単位での概算位置 (左上セルA1を(0,0)とする)
    # これは非常に単純な仮の計算であり、実際のExcelのセルサイズは複雑です
    # 列幅や行高さはopenpyxlで取得可能ですが、ここでは簡略化
    x_pix = (col_idx - 1) * col_width_pix
    y_pix = (row_idx - 1) * row_height_pix

    # ピクセルをインチに変換
    x_inches = x_pix / dpi
    y_inches = y_pix / dpi

    return Inches(x_inches), Inches(y_inches)

def insert_images_to_pptx(pptx_filepath: str, image_folder_path: str, regions_and_coords: list, excel_conv_params: dict, template_filepath: str = None):
    """
    指定された画像フォルダ内の画像を読み込み、その領域をPowerPointスライドの指定座標に貼り付けます。
    テンプレートファイルの1,2ページを保持し、3ページ目を画像の数だけ複製して使用します。
    画像ごとに新しいスライドを作成し、スライド右上に画像ファイル名を表記します。

    Args:
        pptx_filepath (str): 出力するPowerPointファイルのパス。
        image_folder_path (str): 画像が保存されているフォルダのパス。
        regions_and_coords (list): 領域とセル座標のペアのリスト。
                                   例: [{"img_region": [x1, y1, x2, y2], "excel_pos": "B2"}, ...]
                                   img_region: [left, upper, right, lower] (Pillowのcrop形式)
        excel_conv_params (dict): Excelのセル座標をインチに変換するためのパラメータ。
                                  例: {"col_width_pix": 64, "row_height_pix": 20, "dpi": 96}
        template_filepath (str, optional): テンプレートとして使用するPowerPointファイルのパス。
                                           指定しない場合は新規プレゼンテーションを作成します。
    """
    prs = None
    if template_filepath and os.path.exists(template_filepath):
        prs = Presentation(template_filepath)
        print(f"Using template PowerPoint file: {template_filepath}")
    else:
        prs = Presentation()
        print("Creating a new PowerPoint presentation without a template.")

    supported_extensions = ('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff')

    def _duplicate_slide(pres, slide_to_duplicate):
        """指定されたスライドを複製し、プレゼンテーションの最後に追加します。"""
        # スライドのレイアウトを取得
        slide_layout = slide_to_duplicate.slide_layout
        new_slide = pres.slides.add_slide(slide_layout)

        # 複製元のスライドの各シェイプを新しいスライドにコピー
        for shape in slide_to_duplicate.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                print(f"  Warning: Picture shape found in template slide. Skipping direct copy of image data.")
                # 画像の複製は一時ファイル経由で行うため、ここでは何もしない
            elif shape.has_text_frame:
                new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                new_shape.text_frame.text = "" # Clear the text content
                # フォントなどのプロパティもコピー可能だが、ここでは簡略化
            elif shape.is_placeholder:
                # プレースホルダーの場合は、対応するレイアウトプレースホルダーを使用
                placeholder = None
                for p in slide_layout.placeholders:
                    if p.placeholder_format.idx == shape.placeholder_format.idx:
                        placeholder = p
                        break
                if placeholder:
                    new_shape = new_slide.shapes.placeholders[placeholder.placeholder_format.idx]
                    if new_shape.has_text_frame:
                        new_shape.text_frame.text = "" # プレースホルダーの内容をクリア
                else:
                    new_shape = new_slide.shapes.add_shape(shape.shape_type, shape.left, shape.top, shape.width, shape.height)
            elif shape.shape_type != MSO_SHAPE_TYPE.GROUP: # グループ化されたシェイプは個別に処理しない
                 new_shape = new_slide.shapes.add_shape(shape.shape_type, shape.left, shape.top, shape.width, shape.height)

            # TODO: その他のシェイプタイプ（テーブル、グラフなど）は別途ハンドリングが必要
        return new_slide

    image_files = sorted([f for f in os.listdir(image_folder_path) if f.lower().endswith(supported_extensions)])

    if not image_files:
        print("No supported image files found in the specified folder.")
        # 画像がない場合でもテンプレートがあれば保存は試みる
        if prs and (template_filepath and os.path.exists(template_filepath)):
            print("No images to insert, but saving template as is.")
            try:
                prs.save(pptx_filepath)
                print(f"PowerPoint file saved successfully to {pptx_filepath}")
            except Exception as e:
                print(f"Error saving PowerPoint file: {e}")
        return

    # テンプレート使用時のスライド処理ロジック
    if template_filepath and os.path.exists(template_filepath):
        if len(prs.slides) < 3:
            raise ValueError("Template PowerPoint file must have at least 3 slides to use the third slide as a template.")
        
        template_slide = prs.slides[2] # 0-indexed for 3rd slide

        slides_to_process = []
        # 最初の画像は既存の3ページ目を使用
        slide_to_use = prs.slides[2]
        # 既存のコンテンツ（特にテキストボックスや画像）をクリアする
        shapes_to_remove = [s for s in list(slide_to_use.shapes) if s.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER]
        for shape in shapes_to_remove:
            try:
                shape.element.getparent().remove(shape.element)
            except Exception as e:
                print(f"  Warning: Could not remove shape from template slide: {e}")
        
        # プレースホルダーのテキストをクリア
        for placeholder in slide_to_use.placeholders:
            if placeholder.has_text_frame:
                placeholder.text_frame.text = ""
        
        slides_to_process.append(slide_to_use)

        # 残りの画像は3ページ目を複製して使用
        for _ in range(1, len(image_files)):
            slides_to_process.append(_duplicate_slide(prs, template_slide))

    else: # テンプレートファイルが指定されていない、または存在しない場合
        # 新規作成したプレゼンテーションの場合、スライドを最初から追加
        # prsは既に初期化済み
        slides_to_process = []
        blank_slide_layout = prs.slide_layouts[6] # 通常、6番目が空白レイアウト
        for _ in range(len(image_files)):
            slides_to_process.append(prs.slides.add_slide(blank_slide_layout))

    for i, image_filename in enumerate(image_files):
        image_path = os.path.join(image_folder_path, image_filename)
        current_slide = slides_to_process[i]
        
        print(f"Processing image: {image_filename} on slide {prs.slides.index(current_slide) + 1}")

        # スライド右上に画像ファイル名を表記
        left = Inches(prs.slide_width.inches - 2)
        top = Inches(0.1)
        width = Inches(1.9)
        height = Inches(0.5)
        textbox = current_slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = image_filename
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.font.size = Pt(10)

        try:
            original_image = Image.open(image_path)

            for item in regions_and_coords:
                img_region = item["img_region"]
                excel_pos = item["excel_pos"]

                # 画像領域を切り抜き
                cropped_image = original_image.crop(img_region)

                # 切り抜いた画像を一時ファイルとして保存
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                    temp_image_path = temp_file.name
                    cropped_image.save(temp_image_path)

                # Excelセル座標をPowerPointのインチ座標に変換
                x_inches, y_inches = excel_coord_to_inches(excel_pos, excel_conv_params)

                # 画像をスライドに貼り付け
                pic = current_slide.shapes.add_picture(temp_image_path, x_inches, y_inches)
                print(f"  - Cropped region {img_region} from {image_filename} and inserted at {excel_pos} ({x_inches.inches:.2f}in, {y_inches.inches:.2f}in)")

                # 一時ファイルを削除
                os.remove(temp_image_path)

        except Exception as e:
            print(f"Error processing {image_filename}: {e}")
            continue

    # プレゼンテーションを保存
    try:
        prs.save(pptx_filepath)
        print(f"PowerPoint file saved successfully to {pptx_filepath}")
    except Exception as e:
        print(f"Error saving PowerPoint file: {e}")
    finally:
        pass

if __name__ == '__main__':
    current_dir = os.path.dirname(os.path.abspath(__file__))
    config_path = os.path.join(current_dir, "config.json")
    image_dir = os.path.join(current_dir, "img")
    output_pptx_file = os.path.join(current_dir, "output_images.pptx")

    try:
        with open(config_path, 'r', encoding='utf-8') as f:
            config = json.load(f)

        regions_and_coords = config.get("image_regions_and_excel_coords", [])
        excel_conversion_parameters = config.get("excel_to_pptx_conversion_params", {})

        # 関数を実行
        template_file = config.get("pptx_template_file", None)
        template_pptx_path = os.path.join(current_dir, template_file) if template_file else None
        
        insert_images_to_pptx(output_pptx_file, image_dir, regions_and_coords, excel_conversion_parameters, template_pptx_path)

    except FileNotFoundError:
        print(f"Error: config file not found at {config_path}")
    except json.JSONDecodeError:
        print(f"Error: Could not decode JSON from {config_path}. Check file format.")
    except Exception as e:
        print(f"An unexpected error occurred: {e}")


